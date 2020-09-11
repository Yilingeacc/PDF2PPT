# -*-coding: UTF-8 -*-
# Author Shi

from pdf2image import convert_from_path
from pptx import Presentation
import os

cnt = 0
os.mkdir('output')
file_list = os.listdir('convert')
for file in file_list:
    os.mkdir('temp')
    convert_from_path('convert/' + file, output_folder='temp', fmt='JPG')
    img_list = os.listdir('temp')
    for i in range(0, len(img_list)):
        os.rename('temp/' + img_list[i], 'temp/' + str(i) + ".jpg")
    prs = Presentation('template.pptx')

    cover_path = '0.jpg'
    slide = prs.slides[0]
    pic = slide.shapes.add_picture(cover_path, 0, 0)
    SLD_LAYOUT_TITLE_AND_CONTENT = 6
    img_list = os.listdir('temp')
    for i in range(0, len(img_list)):
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'temp/' + str(i) + '.jpg'
        pic = slide.shapes.add_picture(img_path, 0, 0)
        cnt = cnt + 1
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]
    prs.save('output/' + file + '.pptx')
    for img in img_list:
        os.remove('temp/' + img)
    os.rmdir('temp')
print(cnt)
