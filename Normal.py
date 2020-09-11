# -*-coding: UTF-8 -*-
# Author Shi

from pdf2image import convert_from_path
from pptx import Presentation
import cv2
import os

cnt = 0
if not (os.listdir().__contains__('output')):
    os.mkdir('output')
file_list = os.listdir('convert')
for file in file_list:
    os.mkdir('temp')
    convert_from_path('convert/' + file, output_folder='temp', fmt='JPG')
    img_list = os.listdir('temp')
    i = 0
    for img in img_list:
        src = cv2.imread('temp/' + img_list[i])
        src = cv2.resize(src, (1500, 1125), interpolation=cv2.INTER_AREA)
        cv2.imwrite("temp/" + str(i) + ".jpg", src)
        os.remove('temp/' + img_list[i])
        i = i + 1

    prs = Presentation('template.pptx')

    SLD_LAYOUT_TITLE_AND_CONTENT = 6
    img_list = os.listdir('temp')
    for i in range(0, len(img_list)):
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = prs.slides.add_slide(slide_layout)
        img_path = 'temp/' + str(i) + '.jpg'
        pic = slide.shapes.add_picture(img_path, left=0, top=0)
        cnt = cnt + 1
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]
    prs.save('output/' + file[0:-4] + '.pptx')
    for img in img_list:
        os.remove('temp/' + img)
    os.rmdir('temp')
print(cnt)
